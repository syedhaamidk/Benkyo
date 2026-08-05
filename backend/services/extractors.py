"""
services/extractors.py — Format-specific text extraction.

Each public function accepts a file path (str or Path) and returns a list of
chunk-candidates:

    [{"text": str, "source_filename": str, "page_number": int}, ...]

page_number is 1-based for PDFs/PPTX, 0 for plain text / images.

Graceful degradation:
  - If PyMuPDF is not installed, PDF extraction raises ImportError.
  - If pytesseract / Tesseract is not available, OCR falls back to skipping
    scanned pages rather than crashing.
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _tesseract_available() -> bool:
    """Return True if pytesseract + tesseract binary are both usable."""
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


_TESSERACT_OK: bool | None = None  # cached after first check


def _ocr_image(pil_image: Any) -> str:
    """Run OCR on a PIL Image; returns empty string if Tesseract unavailable."""
    global _TESSERACT_OK
    if _TESSERACT_OK is None:
        _TESSERACT_OK = _tesseract_available()
        if not _TESSERACT_OK:
            logger.warning(
                "Tesseract not found — OCR will be skipped for scanned pages/images. "
                "Install from https://github.com/UB-Mannheim/tesseract/wiki"
            )
    if not _TESSERACT_OK:
        return ""
    import pytesseract
    return pytesseract.image_to_string(pil_image)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_pdf(file_path: str | Path) -> list[dict]:
    """
    Extract text from a PDF using PyMuPDF.
    Falls back to Tesseract OCR for pages that yield less than 20 characters
    of text (typically scanned / image-only pages).
    """
    import fitz  # PyMuPDF

    file_path = Path(file_path)
    filename = file_path.name
    results: list[dict] = []

    doc = fitz.open(str(file_path))
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()

        if len(text) < 20:
            # Likely a scanned page — try OCR
            try:
                from PIL import Image
                import io

                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                text = _ocr_image(img).strip()
            except Exception as exc:
                logger.debug("OCR failed for page %d: %s", page_num, exc)
                text = ""

        if text:
            results.append(
                {"text": text, "source_filename": filename, "page_number": page_num}
            )

    doc.close()
    return results


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx(file_path: str | Path) -> list[dict]:
    """Extract slide text + speaker notes from a PowerPoint file."""
    from pptx import Presentation

    file_path = Path(file_path)
    filename = file_path.name
    results: list[dict] = []

    prs = Presentation(str(file_path))
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = "\n".join(
                para.text.strip()
                for para in notes_tf.paragraphs
                if para.text.strip()
            )
            if notes_text:
                parts.append(f"[Speaker notes]: {notes_text}")

        text = "\n".join(parts).strip()
        if text:
            results.append(
                {"text": text, "source_filename": filename, "page_number": slide_num}
            )

    return results


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_docx(file_path: str | Path) -> list[dict]:
    """Extract paragraph text from a Word document."""
    from docx import Document as DocxDocument

    file_path = Path(file_path)
    filename = file_path.name

    doc = DocxDocument(str(file_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)

    if not text:
        return []

    return [{"text": text, "source_filename": filename, "page_number": 0}]


# ---------------------------------------------------------------------------
# Plain text / Markdown
# ---------------------------------------------------------------------------

def extract_txt(file_path: str | Path) -> list[dict]:
    """Read a plain text or Markdown file."""
    file_path = Path(file_path)
    text = file_path.read_text(encoding="utf-8", errors="replace").strip()
    if not text:
        return []
    return [{"text": text, "source_filename": file_path.name, "page_number": 0}]


# ---------------------------------------------------------------------------
# Image (photos of notes, diagrams, etc.)
# ---------------------------------------------------------------------------

def extract_image(file_path: str | Path) -> list[dict]:
    """Run Tesseract OCR directly on an image file."""
    from PIL import Image

    file_path = Path(file_path)
    img = Image.open(str(file_path))
    text = _ocr_image(img).strip()
    if not text:
        return []
    return [{"text": text, "source_filename": file_path.name, "page_number": 0}]


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXTENSION_MAP = {
    ".pdf": extract_pdf,
    # NOTE: .ppt / .doc (legacy binary Office formats) are intentionally NOT
    # mapped here. python-pptx and python-docx only read the modern zip-based
    # OOXML formats (.pptx/.docx) and raise an unhandled exception on legacy
    # binary files — mapping them here would make uploads silently fail.
    # Ask the user to re-save as .pptx/.docx, or add a dedicated legacy
    # parser (e.g. antiword/textract) if real support is needed.
    ".pptx": extract_pptx,
    ".docx": extract_docx,
    ".txt": extract_txt,
    ".md": extract_txt,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".webp": extract_image,
    ".bmp": extract_image,
    ".tiff": extract_image,
    ".tif": extract_image,
}


def extract_file(file_path: str | Path) -> list[dict]:
    """
    Auto-dispatch to the correct extractor based on file extension.
    Raises ValueError for unsupported types.
    """
    ext = Path(file_path).suffix.lower()
    extractor = _EXTENSION_MAP.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext!r}")
    return extractor(file_path)
